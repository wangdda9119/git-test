"""
Render 1.html .. 15.html at 1280x720 and build a widescreen PPTX (one image per slide).
"""
from __future__ import annotations

import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

try:
    from playwright.sync_api import sync_playwright
except ImportError:
    print("Install dependencies: pip install -r requirements.txt && playwright install chromium")
    sys.exit(1)

ROOT = Path(__file__).resolve().parent
SLIDE_W, SLIDE_H = 1280, 720
# PowerPoint widescreen 16:9 (inches)
PPTX_W = Inches(13.333333)
PPTX_H = Inches(7.5)


def html_files() -> list[Path]:
    files = []
    for i in range(1, 16):
        p = ROOT / f"{i}.html"
        if not p.is_file():
            raise FileNotFoundError(f"Missing: {p}")
        if p.stat().st_size < 400:
            print(
                f"경고: {p.name} 크기가 매우 작습니다({p.stat().st_size}B). "
                "에디터에서 내용을 채운 뒤 디스크에 저장(Ctrl+K S 등)했는지 확인하세요.",
                file=sys.stderr,
            )
        files.append(p)
    return files


def render_slides(paths: list[Path], out_dir: Path) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    pngs: list[Path] = []
    with sync_playwright() as p:
        browser = None
        last_err: Exception | None = None
        for channel in ("msedge", "chrome", None):
            try:
                if channel:
                    browser = p.chromium.launch(headless=True, channel=channel)
                else:
                    browser = p.chromium.launch(headless=True)
                break
            except Exception as e:
                last_err = e
        if browser is None:
            raise RuntimeError(
                "브라우저를 시작할 수 없습니다. Edge 또는 Chrome 설치 후 다시 시도하세요."
            ) from last_err
        context = browser.new_context(
            viewport={"width": SLIDE_W, "height": SLIDE_H},
            device_scale_factor=2,
        )
        page = context.new_page()
        for idx, html_path in enumerate(paths, start=1):
            uri = html_path.resolve().as_uri()
            page.goto(uri, wait_until="networkidle", timeout=120_000)
            # Let webfonts settle
            page.wait_for_timeout(500)
            png = out_dir / f"slide_{idx:02d}.png"
            page.screenshot(path=str(png), clip={"x": 0, "y": 0, "width": SLIDE_W, "height": SLIDE_H})
            pngs.append(png)
        browser.close()
    return pngs


def build_pptx(png_paths: list[Path], output: Path) -> None:
    prs = Presentation()
    prs.slide_width = PPTX_W
    prs.slide_height = PPTX_H
    blank = prs.slide_layouts[6]  # blank

    for png in png_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

    prs.save(str(output))


def main() -> None:
    paths = html_files()
    out_pptx = ROOT / "presentation.pptx"
    with tempfile.TemporaryDirectory(prefix="ge_ppt_") as tmp:
        tmp_path = Path(tmp)
        pngs = render_slides(paths, tmp_path)
        build_pptx(pngs, out_pptx)
    print(f"Wrote: {out_pptx}")


if __name__ == "__main__":
    main()
